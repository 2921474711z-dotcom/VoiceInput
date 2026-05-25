from __future__ import annotations

import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[2]
VENDOR = ROOT / ".vendor"
if VENDOR.exists():
    sys.path.insert(0, str(VENDOR))

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_DIR = ROOT / "docs" / "presentation"
ASSET_DIR = ROOT / ".tmp" / "pptx-assets"
OUT_FILE = OUT_DIR / "VoiceInput-Pro-平台介绍与验收演示稿.pptx"

SCREENSHOTS = {
    "workbench": ASSET_DIR / "01-workbench.png",
    "history": ASSET_DIR / "02-history.png",
    "exports": ASSET_DIR / "03-exports.png",
    "hotwords": ASSET_DIR / "04-hotwords.png",
    "config": ASSET_DIR / "05-config.png",
}

NAVY = RGBColor(24, 59, 91)
NAVY_DARK = RGBColor(15, 34, 51)
TEAL = RGBColor(14, 116, 144)
TEAL_SOFT = RGBColor(224, 242, 247)
GOLD = RGBColor(244, 183, 64)
SLATE = RGBColor(30, 41, 59)
MUTED = RGBColor(91, 107, 122)
LIGHT = RGBColor(247, 250, 252)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(217, 226, 236)


def add_textbox(slide, left, top, width, height, text, *, font_size=18, color=SLATE,
                bold=False, font_name="Microsoft YaHei", align=PP_ALIGN.LEFT,
                valign=MSO_ANCHOR.TOP, margin=0.06, line_spacing=1.18):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, left, top, width, height, items, *, font_size=20, color=SLATE, level0=0):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.level = level0
        p.bullet = True
        p.space_after = Pt(8)
        p.line_spacing = 1.18
    return box


def add_card(slide, left, top, width, height, *, fill=WHITE, line=LINE,
             shape_type=MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1)
    return shape


def add_section_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.55), Inches(0.35), Inches(8.2), Inches(0.5), title,
                font_size=28, bold=True, color=NAVY, margin=0)
    if subtitle:
        add_textbox(slide, Inches(0.58), Inches(0.82), Inches(8.7), Inches(0.42), subtitle,
                    font_size=11.5, color=MUTED, margin=0)
    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.38), Inches(0.36), Inches(0.08), Inches(0.72))
    accent.fill.solid()
    accent.fill.fore_color.rgb = TEAL
    accent.line.fill.background()


def add_image_fit(slide, image_path: Path, left, top, width, height):
    with Image.open(image_path) as img:
        img_ratio = img.width / img.height
    box_ratio = width / height
    if img_ratio >= box_ratio:
        final_w = width
        final_h = width / img_ratio
        final_left = left
        final_top = top + (height - final_h) / 2
    else:
        final_h = height
        final_w = height * img_ratio
        final_top = top
        final_left = left + (width - final_w) / 2
    slide.shapes.add_picture(str(image_path), final_left, final_top, width=final_w, height=final_h)


def add_page_mock(slide, image_path: Path, left, top, width, height, caption):
    add_card(slide, left, top, width, height)
    img_h = height - Inches(0.5)
    add_image_fit(slide, image_path, left + Inches(0.08), top + Inches(0.08), width - Inches(0.16), img_h - Inches(0.14))
    add_textbox(slide, left + Inches(0.14), top + img_h, width - Inches(0.28), Inches(0.28), caption,
                font_size=10.5, color=MUTED, margin=0, valign=MSO_ANCHOR.MIDDLE)


def add_metric(slide, left, top, width, height, value, label, note):
    add_card(slide, left, top, width, height, fill=WHITE)
    add_textbox(slide, left + Inches(0.18), top + Inches(0.18), width - Inches(0.36), Inches(0.48), value,
                font_size=26, bold=True, color=NAVY, margin=0)
    add_textbox(slide, left + Inches(0.18), top + Inches(0.66), width - Inches(0.36), Inches(0.3), label,
                font_size=11, bold=True, color=TEAL, margin=0)
    add_textbox(slide, left + Inches(0.18), top + Inches(1.0), width - Inches(0.36), height - Inches(1.14), note,
                font_size=10.5, color=MUTED, margin=0)


def draw_flow(slide, steps):
    start_x = 0.55
    y = 1.75
    w = 1.45
    h = 0.88
    gap = 0.18
    for idx, step in enumerate(steps):
        x = Inches(start_x + idx * (w + gap))
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, Inches(y), Inches(w), Inches(h))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE if idx % 2 == 0 else TEAL_SOFT
        box.line.color.rgb = TEAL
        box.line.width = Pt(1.2)
        add_textbox(slide, x + Inches(0.1), Inches(y) + Inches(0.12), Inches(w - 0.2), Inches(h - 0.24), step,
                    font_size=12.5, bold=True, color=SLATE, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, margin=0)
        if idx < len(steps) - 1:
            line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x + Inches(w), Inches(y + h / 2), x + Inches(w + gap), Inches(y + h / 2))
            line.line.color.rgb = TEAL
            line.line.width = Pt(2)


def build_presentation():
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.author = "OpenAI Codex"
    prs.core_properties.title = "VoiceInput Pro 平台介绍与验收演示稿"
    prs.core_properties.subject = "VoiceInput Pro 产品演示"
    prs.core_properties.comments = "基于当前本地站点页面与功能整理"

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY_DARK
    ribbon = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.72), Inches(0.86), Inches(1.18), Inches(0.34))
    ribbon.fill.solid()
    ribbon.fill.fore_color.rgb = GOLD
    ribbon.line.fill.background()
    add_textbox(slide, Inches(0.78), Inches(0.9), Inches(1.05), Inches(0.24), "平台演示稿",
                font_size=10.5, bold=True, color=NAVY_DARK, margin=0, valign=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, Inches(0.72), Inches(1.45), Inches(6.7), Inches(1.2), "VoiceInput Pro",
                font_size=28, bold=True, color=WHITE, margin=0)
    add_textbox(slide, Inches(0.72), Inches(2.22), Inches(7.8), Inches(1.0),
                "一套面向办公与技术场景的智能语音输入工作台\n从音频上传、识别优化，到历史沉淀、导出交付，页面与数据链路完整打通。",
                font_size=18, color=WHITE, margin=0, line_spacing=1.28)
    add_textbox(slide, Inches(0.72), Inches(6.48), Inches(4.2), Inches(0.34),
                "适合课堂答辩、项目验收、功能演示", font_size=11.5, color=RGBColor(206, 216, 226), margin=0)
    add_page_mock(slide, SCREENSHOTS["workbench"], Inches(8.0), Inches(1.02), Inches(4.62), Inches(5.75), "工作台首页实拍")

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT
    add_section_title(slide, "这套平台解决什么问题", "它不是简单的“语音转文字”，而是把转写、整理、复核和交付串成一条完整工作流。")
    add_metric(slide, Inches(0.65), Inches(1.45), Inches(2.9), Inches(1.65), "6 类场景", "按场景组织输出", "支持会议纪要、工作汇报、正式表达、Markdown 笔记、代码注释、聊天回复。")
    add_metric(slide, Inches(3.72), Inches(1.45), Inches(2.9), Inches(1.65), "双文本输出", "原文 + 优化文并行展示", "用户能直接比较识别结果和整理结果，知道系统改了什么。")
    add_metric(slide, Inches(6.79), Inches(1.45), Inches(2.9), Inches(1.65), "闭环管理", "处理后还能继续用", "结果可以进入历史、参与人工校对、导出成文档，也能进入统计分析。")
    add_metric(slide, Inches(9.86), Inches(1.45), Inches(2.82), Inches(1.65), "面向演示", "每个入口都有实际动作", "页面保留的按钮、筛选、导出、统计，都对应后端能力，而不是静态摆设。")
    add_bullets(slide, Inches(0.82), Inches(3.55), Inches(6.0), Inches(2.6), [
        "对普通转写工具来说，输出往往停留在“一段文字”；对这套平台来说，重点是把文字整理成可直接使用的结果。",
        "平台适合老师、评审或团队现场查看，因为能清楚看到每一步是怎么走的，数据落到了哪里，后续还能怎么复用。",
        "如果把它放到日常办公场景里，最容易落地的用法就是会议纪要、口头汇报整理和技术内容录入。"
    ], font_size=16)
    add_card(slide, Inches(7.15), Inches(3.42), Inches(5.2), Inches(2.85), fill=WHITE)
    add_textbox(slide, Inches(7.38), Inches(3.66), Inches(4.7), Inches(0.36), "一句话理解", font_size=16, bold=True, color=TEAL, margin=0)
    add_textbox(slide, Inches(7.38), Inches(4.15), Inches(4.55), Inches(1.45),
                "它更像一张“语音内容工作台”，不是录完就结束，而是让结果继续流进历史、热词、配置和导出体系里。",
                font_size=19, bold=True, color=SLATE, margin=0, line_spacing=1.25)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT
    add_section_title(slide, "页面结构一眼看清", "顶部导航把核心能力拆成五块，符合从“生产结果”到“管理结果”的使用顺序。")
    nav_cards = [
        ("工作台", "选择场景、上传音频、开始处理、查看原文和优化文。"),
        ("历史记录", "集中查看处理结果，支持筛选、详情、再次优化和删除。"),
        ("导出中心", "统一管理 DOCX、Markdown、TXT、JSON 等导出文件。"),
        ("热词管理", "维护术语和误识别修正规则，让后续处理更贴近真实表达。"),
        ("模型配置", "管理模板、模型路线、统计指标和默认处理策略。"),
    ]
    for idx, (title, desc) in enumerate(nav_cards):
        row = idx // 3
        col = idx % 3
        left = Inches(0.75 + col * 4.1)
        top = Inches(1.55 + row * 2.2)
        width = Inches(3.72 if row == 0 else 5.82 if col == 0 else 5.65)
        if row == 1 and col == 1:
            left = Inches(6.67)
        add_card(slide, left, top, width, Inches(1.65), fill=WHITE)
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left + Inches(0.18), top + Inches(0.2), Inches(0.38), Inches(0.38))
        dot.fill.solid()
        dot.fill.fore_color.rgb = TEAL
        dot.line.fill.background()
        add_textbox(slide, left + Inches(0.68), top + Inches(0.16), width - Inches(0.88), Inches(0.28), title,
                    font_size=17, bold=True, color=NAVY, margin=0)
        add_textbox(slide, left + Inches(0.22), top + Inches(0.62), width - Inches(0.44), Inches(0.76), desc,
                    font_size=12, color=MUTED, margin=0)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT
    add_section_title(slide, "工作台：从一段音频到一份可用文本", "这是平台的主入口，用户会在这里完成上传、处理、查看结果和保存历史。")
    add_page_mock(slide, SCREENSHOTS["workbench"], Inches(0.62), Inches(1.45), Inches(7.25), Inches(5.3), "工作台页面")
    add_bullets(slide, Inches(8.18), Inches(1.62), Inches(4.55), Inches(4.8), [
        "左侧先选场景，系统会根据场景改变输出方式，比如会议纪要更强调议题、结论和待办，代码注释则更看重术语保留。",
        "中间完成音频上传后，就能走完整处理链路：识别、热词修正、文本优化，再把结果和统计信息写回。",
        "结果区同时展示原始识别文本和优化后文本，方便对比，也方便后续人工校对。",
        "同一页还能继续做“保存到历史”“再次处理”“人工校对”“生成导出”等动作，减少来回切换。"
    ], font_size=15)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT
    add_section_title(slide, "历史记录与导出中心", "处理结果不会停留在当前会话里，平台把“查找”“复用”“交付”放到了同一套体系。")
    add_page_mock(slide, SCREENSHOTS["history"], Inches(0.62), Inches(1.55), Inches(6.2), Inches(4.95), "历史记录页面")
    add_page_mock(slide, SCREENSHOTS["exports"], Inches(6.98), Inches(1.55), Inches(5.72), Inches(4.95), "导出中心页面")
    add_bullets(slide, Inches(0.8), Inches(6.62), Inches(12.0), Inches(0.55), [
        "历史页支持关键词搜索、场景筛选、详情查看、人工校对、再次优化和导出；导出中心则把文件生成记录统一沉淀下来，方便交付与追溯。"
    ], font_size=13.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT
    add_section_title(slide, "热词管理：把专业表达纠正到位", "热词模块不是摆设，它直接影响后续处理结果，尤其适合技术词、缩写词和常见误识别场景。")
    add_page_mock(slide, SCREENSHOTS["hotwords"], Inches(0.62), Inches(1.48), Inches(7.08), Inches(5.28), "热词管理页面")
    add_card(slide, Inches(8.0), Inches(1.62), Inches(4.5), Inches(4.65), fill=WHITE)
    add_textbox(slide, Inches(8.22), Inches(1.88), Inches(3.8), Inches(0.32), "这个模块的价值", font_size=16, bold=True, color=TEAL, margin=0)
    add_bullets(slide, Inches(8.22), Inches(2.32), Inches(3.9), Inches(3.5), [
        "支持分类维护，页面里能看到技术术语、产品名称、英文缩写、自定义词库等类别。",
        "支持新增、编辑、删除、导入和导出，适合快速沉淀一套自己的行业词表。",
        "热词启用后会直接参与文本修正，不是只存着看，后面的统计里也能看到热词命中情况。",
        "对答辩或验收来说，这一页很容易体现项目差异化，因为它说明平台不是“通用转写完事”，而是能针对场景做修正。"
    ], font_size=14.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT
    add_section_title(slide, "模型配置与使用统计", "平台不仅能跑，还允许调策略、存模板、看统计，这让它更像一套可运维的产品。")
    add_page_mock(slide, SCREENSHOTS["config"], Inches(0.62), Inches(1.48), Inches(7.08), Inches(5.28), "模型配置页面")
    add_bullets(slide, Inches(8.0), Inches(1.72), Inches(4.55), Inches(4.95), [
        "配置页可以设置识别模型、优化模型、语言、领域、输出格式、热词开关和成本/质量取向。",
        "模板能力让常用策略能被保存下来，工作台在创建任务时直接复用，减少重复设置。",
        "统计部分会回看任务耗时、平均字数、热词命中率、近 7 天趋势和场景分布，适合拿来说明平台真实使用情况。",
        "如果从产品完整度看，这一块很关键，因为它把“用起来”延伸到了“长期维护和复盘”。"
    ], font_size=14.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT
    add_section_title(slide, "核心处理链路", "从后端实现看，这套平台并不是前端本地拼结果，而是真正经过存储、识别、优化和落库。")
    draw_flow(slide, ["上传音频", "对象存储", "语音识别", "热词修正", "文本优化", "结果回写", "历史/导出"])
    add_bullets(slide, Inches(0.8), Inches(3.15), Inches(12.0), Inches(3.1), [
        "音频先进入存储层，再由后端任务服务读取，避免文件只停留在浏览器内存里。",
        "识别阶段拿到原始文本后，会按当前场景和热词规则继续修正，再把文本交给优化模型做结构化整理。",
        "任务完成后，会把标题、摘要、原文、优化文、Markdown、耗时、热词命中和成本估算一起写回数据库。",
        "这样一来，工作台、历史记录、导出中心、统计页看到的是同一份结果，不会出现各页数据断层。"
    ], font_size=15)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT
    add_section_title(slide, "技术架构与部署方式", "这套平台采用前后端分离设计，适合本地演示，也保留了后续扩展空间。")
    layer_titles = [
        ("前端层", "React + Vite + TypeScript\n负责工作台、历史、热词、导出和配置等页面。"),
        ("后端层", "Spring Boot 3 + JPA + Flyway\n负责任务处理、接口编排、导出生成和统计聚合。"),
        ("数据与中间件", "PostgreSQL + Redis + MinIO\n分别承担持久化、状态/缓存与文件对象存储。"),
        ("部署方式", "Docker Compose 一键拉起\n适合课堂环境和演示机器快速复现。"),
    ]
    for idx, (title, desc) in enumerate(layer_titles):
        left = Inches(0.82 + idx * 3.15)
        add_card(slide, left, Inches(1.9), Inches(2.75), Inches(3.6), fill=WHITE)
        tag = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, Inches(1.9), Inches(2.75), Inches(0.24))
        tag.fill.solid()
        tag.fill.fore_color.rgb = TEAL if idx % 2 == 0 else NAVY
        tag.line.fill.background()
        add_textbox(slide, left + Inches(0.18), Inches(2.28), Inches(2.35), Inches(0.34), title,
                    font_size=16, bold=True, color=SLATE, margin=0)
        add_textbox(slide, left + Inches(0.18), Inches(2.7), Inches(2.35), Inches(2.32), desc,
                    font_size=13, color=MUTED, margin=0, line_spacing=1.28)
    add_textbox(slide, Inches(0.86), Inches(6.05), Inches(11.8), Inches(0.5),
                "项目当前支持本地通过 Docker 启动核心服务，适合在验收现场直接演示页面、接口和导出效果。",
                font_size=13.5, color=MUTED, margin=0)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT
    add_section_title(slide, "为什么这套平台适合拿来验收", "评审最关心的通常不是“页面多不多”，而是“功能是不是真的跑得通”。")
    items = [
        ("入口真实可用", "上传、处理、保存历史、导出、热词维护、配置保存和统计查看，都对应实际行为。"),
        ("数据链路完整", "任务结果不是页面假数据，而是能够回到历史、导出中心和统计页继续使用。"),
        ("场景化差异明确", "六类输入场景让结果不是一刀切输出，更容易展示产品设计思路。"),
        ("便于现场演示", "从工作台开始，一步步能走到历史、热词、配置和导出，流程清晰，讲起来也顺。"),
    ]
    for idx, (title, desc) in enumerate(items):
        left = Inches(0.82 + (idx % 2) * 6.05)
        top = Inches(1.6 + (idx // 2) * 2.0)
        add_card(slide, left, top, Inches(5.62), Inches(1.58), fill=WHITE)
        badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, left + Inches(0.18), top + Inches(0.22), Inches(0.42), Inches(0.42))
        badge.fill.solid()
        badge.fill.fore_color.rgb = GOLD
        badge.line.fill.background()
        add_textbox(slide, left + Inches(0.75), top + Inches(0.17), Inches(4.3), Inches(0.3), title,
                    font_size=15.5, bold=True, color=NAVY, margin=0)
        add_textbox(slide, left + Inches(0.22), top + Inches(0.63), Inches(5.05), Inches(0.68), desc,
                    font_size=12.5, color=MUTED, margin=0)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT
    add_section_title(slide, "建议的现场演示路径", "如果是答辩或平台验收，按这条线讲，节奏最顺，也最容易把亮点讲明白。")
    demo_steps = [
        "1. 在工作台选择一个输入场景，上传音频，展示原文和优化文的对比。",
        "2. 把结果保存到历史记录，再进入历史页说明结果如何被沉淀下来。",
        "3. 打开热词管理，补一条术语修正规则，说明平台为什么适合技术和专业场景。",
        "4. 回到工作台或历史详情，再次处理一次内容，对比修正前后的变化。",
        "5. 最后进入导出中心和模型配置页，展示这套平台不仅能产出结果，还能形成交付文件和统计结果。",
    ]
    add_bullets(slide, Inches(0.85), Inches(1.6), Inches(7.2), Inches(4.8), demo_steps, font_size=17)
    add_card(slide, Inches(8.35), Inches(1.75), Inches(4.0), Inches(3.7), fill=WHITE)
    add_textbox(slide, Inches(8.58), Inches(2.0), Inches(3.3), Inches(0.3), "讲解时建议强调", font_size=16, bold=True, color=TEAL, margin=0)
    add_bullets(slide, Inches(8.58), Inches(2.42), Inches(3.1), Inches(2.6), [
        "页面不是静态稿",
        "结果能回流到后续模块",
        "热词和配置会影响下一次处理",
        "导出文件可以直接交付"
    ], font_size=13.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY_DARK
    add_textbox(slide, Inches(0.78), Inches(1.1), Inches(5.8), Inches(0.7), "总结", font_size=18, bold=True, color=GOLD, margin=0)
    add_textbox(slide, Inches(0.78), Inches(1.75), Inches(7.4), Inches(1.55),
                "VoiceInput Pro 已经具备一套完整的产品骨架：\n能处理、能复核、能沉淀、能导出，也能解释这套结果是怎么来的。",
                font_size=24, bold=True, color=WHITE, margin=0, line_spacing=1.28)
    add_textbox(slide, Inches(0.78), Inches(3.62), Inches(6.8), Inches(1.6),
                "从项目展示角度看，它的优势不只是页面整齐，而是流程完整、数据成链、讲解有抓手。\n如果后续继续扩展，这套平台也有空间往实时录音、团队协作和更多模型策略上走。",
                font_size=16.5, color=RGBColor(219, 229, 236), margin=0, line_spacing=1.28)
    add_page_mock(slide, SCREENSHOTS["workbench"], Inches(8.12), Inches(1.12), Inches(4.38), Inches(5.48), "平台首页")

    prs.save(str(OUT_FILE))


if __name__ == "__main__":
    build_presentation()
    print(OUT_FILE)
